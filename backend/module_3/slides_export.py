from __future__ import annotations

import base64
import io
import uuid
from typing import Any, Dict, Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from backend.module_3.slides_generation import safe_presentation_name


PRESENTATIONS: Dict[str, Dict[str, Any]] = {}
NAVY = RGBColor(15, 23, 42)
BLUE = RGBColor(37, 99, 235)
SKY = RGBColor(224, 242, 254)
SLATE = RGBColor(71, 85, 105)
WHITE = RGBColor(255, 255, 255)
AMBER = RGBColor(245, 158, 11)


def _background(slide: Any, color: RGBColor = WHITE) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _box(slide: Any, left: float, top: float, width: float, height: float,
         color: RGBColor, radius: bool = False) -> Any:
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _text(slide: Any, text: str, left: float, top: float, width: float, height: float,
          size: int, color: RGBColor = NAVY, bold: bool = False,
          align: PP_ALIGN = PP_ALIGN.LEFT) -> Any:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.font.name = "Aptos"
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    return box


def _bullets(slide: Any, bullets: Iterable[str], left: float, top: float,
             width: float, height: float, size: int = 20) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"• {bullet}"
        paragraph.level = 0
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = NAVY
        paragraph.space_after = Pt(10)


def _footer(slide: Any, number: int, pages: str) -> None:
    _box(slide, 0, 7.34, 13.333, 0.16, BLUE)
    _text(slide, f"{number:02d}", 12.45, 6.92, 0.45, 0.25, 10, SLATE, True, PP_ALIGN.RIGHT)
    if pages:
        _text(slide, f"Source: PDF page(s) {pages}", 0.55, 6.92, 5.5, 0.25, 9, SLATE)


def _visual_panel(slide: Any, recommendation: str, left: float = 8.75,
                  top: float = 1.65, width: float = 3.9, height: float = 4.65,
                  image_base64: str = "") -> None:
    _box(slide, left, top, width, height, SKY, True)
    if image_base64:
        slide.shapes.add_picture(
            io.BytesIO(base64.b64decode(image_base64)),
            Inches(left + 0.18),
            Inches(top + 0.18),
            width=Inches(width - 0.36),
            height=Inches((width - 0.36) / 1.5),
        )
        image_height = (width - 0.36) / 1.5
        _text(slide, "VISUAL DIRECTION", left + 0.25, top + image_height + 0.48,
              width - 0.5, 0.3, 10, BLUE, True)
        _text(slide, recommendation, left + 0.25, top + image_height + 0.82,
              width - 0.5, height - image_height - 1.0, 11, NAVY)
        return
    _box(slide, left + 0.25, top + 0.28, 0.5, 0.08, AMBER, True)
    _text(slide, "VISUAL DIRECTION", left + 0.25, top + 0.47, width - 0.5, 0.3, 11, BLUE, True)
    _text(slide, recommendation or "Use a simple concept diagram based on the key idea.",
          left + 0.25, top + 0.92, width - 0.5, height - 1.2, 17, NAVY)


def _render_slide(prs: Presentation, item: Dict[str, Any], number: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    layout = item.get("layout_recommendation", "content")
    title = str(item.get("title") or f"Slide {number}")
    subtitle = str(item.get("subtitle") or "")
    bullets = item.get("bullets") or []
    visual = str(item.get("picture_recommendation") or "")
    generated_image = str(item.get("generated_image_base64") or "")

    if layout == "title" or number == 1:
        _background(slide, NAVY)
        _box(slide, 0.62, 1.1, 0.12, 4.9, AMBER, True)
        title_width = 7.1 if generated_image else 10.9
        _text(slide, title, 1.05, 1.25, title_width, 2.1, 34, WHITE, True)
        if subtitle:
            _text(slide, subtitle, 1.08, 3.45, title_width, 1.15, 20, RGBColor(191, 219, 254))
        if generated_image:
            _visual_panel(slide, visual, 8.75, 1.25, 3.9, 4.75, generated_image)
        _text(slide, "AI SLIDES GENERATOR", 1.08, 5.72, 3.5, 0.35, 11, RGBColor(147, 197, 253), True)
        return

    _background(slide)
    _text(slide, title, 0.65, 0.38, 11.8, 0.72, 28, NAVY, True)
    _box(slide, 0.67, 1.19, 1.05, 0.07, AMBER, True)
    if subtitle:
        _text(slide, subtitle, 0.68, 1.3, 11.5, 0.45, 14, SLATE)

    if generated_image:
        _bullets(slide, bullets, 0.75, 1.85, 7.55, 4.85, 20)
        _visual_panel(slide, visual, image_base64=generated_image)
    elif layout == "section":
        _box(slide, 0.65, 2.0, 12.0, 3.9, NAVY, True)
        _text(slide, title, 1.25, 2.55, 10.8, 1.2, 30, WHITE, True, PP_ALIGN.CENTER)
        if bullets:
            _text(slide, str(bullets[0]), 1.55, 3.85, 10.2, 0.85, 19,
                  RGBColor(191, 219, 254), False, PP_ALIGN.CENTER)
    elif layout == "quote":
        quote = str(bullets[0]) if bullets else subtitle
        _box(slide, 0.85, 2.0, 11.55, 3.7, SKY, True)
        _text(slide, f'“{quote}”', 1.45, 2.45, 10.35, 2.35, 27, NAVY, True, PP_ALIGN.CENTER)
    elif layout == "two_column" and len(bullets) > 1:
        middle = (len(bullets) + 1) // 2
        _bullets(slide, bullets[:middle], 0.8, 1.9, 5.6, 4.7, 19)
        _box(slide, 6.62, 2.0, 0.03, 4.2, RGBColor(203, 213, 225))
        _bullets(slide, bullets[middle:], 6.9, 1.9, 5.5, 4.7, 19)
    elif visual:
        _bullets(slide, bullets, 0.75, 1.85, 7.55, 4.85, 20)
        _visual_panel(slide, visual)
    else:
        _bullets(slide, bullets, 0.85, 1.9, 11.45, 4.7, 21)
    _footer(slide, number, str(item.get("source_pages") or ""))


def _presentation_bytes(state: Dict[str, Any]) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for index, slide_data in enumerate(state["deck"]["slides"], start=1):
        _render_slide(prs, slide_data, index)
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


def export_slides(state: Dict[str, Any]) -> Dict[str, Any]:
    presentation_id = uuid.uuid4().hex
    filename = f"{safe_presentation_name(state['book']['name'])}-slides.pptx"
    PRESENTATIONS[presentation_id] = {
        "content": _presentation_bytes(state),
        "filename": filename,
    }
    result = {
        "presentation_id": presentation_id,
        "filename": filename,
        "download": f"/slides/{presentation_id}/download",
    }
    return {"presentation_id": presentation_id, "result": result}


def get_presentation(presentation_id: str) -> tuple[bytes, str]:
    presentation = PRESENTATIONS.get(presentation_id)
    if not presentation:
        raise KeyError("Presentation not found or expired.")
    return presentation["content"], presentation["filename"]
